#!/usr/bin/env python3

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BG = RGBColor(0x0D, 0x0D, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0xCC, 0xCC, 0xCC)
QUIET = RGBColor(0x8A, 0x92, 0xA6)
CYAN = RGBColor(0x00, 0xF0, 0xFF)
PURPLE = RGBColor(0xD4, 0xA0, 0xFF)
ORANGE = RGBColor(0xFF, 0x9A, 0x00)
GREEN = RGBColor(0x00, 0xFF, 0x88)
RED = RGBColor(0xFF, 0x44, 0x44)
CARD_FILL = RGBColor(0x17, 0x1A, 0x22)
CARD_FILL_2 = RGBColor(0x12, 0x14, 0x1B)
CARD_LINE = RGBColor(0x74, 0x7B, 0x8E)


def px(value):
    return Inches(value / 144.0)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_orb(slide, x, y, w, h, color, transparency):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, px(x), px(y), px(w), px(h))
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    fill.transparency = transparency
    return shape


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    size,
    color=WHITE,
    bold=False,
    align=PP_ALIGN.LEFT,
    font="Microsoft YaHei",
    vertical=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = vertical
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    p.alignment = align
    return box


def add_card(slide, x, y, w, h, line_color=CARD_LINE, fill_color=CARD_FILL, transparency=0.08, radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    card = slide.shapes.add_shape(radius_shape, px(x), px(y), px(w), px(h))
    fill = card.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    fill.transparency = transparency
    line = card.line
    line.color.rgb = line_color
    line.transparency = 0.4
    line.width = Pt(1)
    return card


def add_footer(slide, text):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, px(50), px(1036), px(1820), px(18))
    bar.line.fill.background()
    fill = bar.fill
    fill.solid()
    fill.fore_color.rgb = CARD_FILL
    fill.transparency = 0.12
    add_textbox(slide, 64, 1020, 700, 30, text, 11, QUIET, False)


def add_header(slide, left_text, right_text="openclaw / mimo-v2-omni"):
    add_textbox(slide, 50, 48, 1000, 40, left_text, 22, WHITE, True)
    add_textbox(slide, 1200, 48, 670, 36, right_text, 14, QUIET, False, align=PP_ALIGN.RIGHT)


def add_chip(slide, x, y, w, h, text, color):
    chip = add_card(slide, x, y, w, h, line_color=color, fill_color=CARD_FILL_2, transparency=0.02)
    chip.line.transparency = 0.15
    add_textbox(slide, x, y + 6, w, h - 6, text, 12, color, True, PP_ALIGN.CENTER)
    return chip


def add_logic_card(slide, x, y, w, h, text, accent=CYAN):
    """Adds a small-form minimal logic card for Coding Plan nodes."""
    add_card(slide, x, y, w, h, line_color=accent, fill_color=CARD_FILL_2, transparency=0.02)
    add_textbox(slide, x + 10, y + 10, w - 20, h - 20, text, 14, WHITE, False, PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE)


def add_comparison_card(slide, x, y, w, h, title, accent, items):
    add_card(slide, x, y, w, h, line_color=accent, fill_color=CARD_FILL, transparency=0.04)
    add_textbox(slide, x + 28, y + 26, w - 56, 38, title, 24, accent, True)
    top = y + 92
    for item in items:
        add_card(slide, x + 28, top, w - 56, 92, line_color=CARD_LINE, fill_color=CARD_FILL_2, transparency=0.02)
        add_textbox(slide, x + 52, top + 24, w - 104, 44, item, 18, WHITE, False)
        top += 112


def add_bar(slide, x, y, w, h, label, value_text, color, muted=False):
    fill_color = CARD_FILL_2 if muted else color
    trans = 0.35 if muted else 0.08
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, px(x), px(y), px(w), px(h))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = fill_color
    bar.fill.transparency = trans
    add_textbox(slide, x - 6, y - 28, 220, 24, label, 12, QUIET if muted else TEXT, False)
    add_textbox(slide, x + w + 10, y - 6, 90, 24, value_text, 12, color if not muted else QUIET, True)


def add_code_window(slide, x, y, w, h):
    add_card(slide, x, y, w, h, line_color=CYAN, fill_color=CARD_FILL_2, transparency=0.01)
    add_textbox(slide, x + 26, y + 16, w - 52, 22, "terminal", 12, QUIET, False, font="Consolas")
    code = (
        "$ python3 run-chain.py plan video-parse-from-url.yaml \\\n"
        "    --var URL='https://v.douyin.com/...' \\\n"
        "    --var WORKDIR='./output' \\\n"
        "    --var PROMPT='analyze structure'\n\n"
        "$ node openrouter-mimo-omni.js \\\n"
        "    --video source-video.mp4 \\\n"
        "    --report analysis.md\n\n"
        "# Keep the command area short.\n"
        "# Show only the line that proves the workflow."
    )
    add_textbox(slide, x + 30, y + 60, w - 60, h - 90, code, 16, WHITE, False, font="Consolas")


def add_repo_card(slide, x, y, w, h, repo_url):
    add_card(slide, x, y, w, h, line_color=CYAN, fill_color=CARD_FILL, transparency=0.03)
    add_textbox(slide, x + 28, y + 24, w - 56, 34, "GitHub / Next Step", 18, CYAN, True)
    add_textbox(slide, x + 28, y + 70, w - 56, 60, repo_url, 18, WHITE, False)
    add_textbox(slide, x + 28, y + 142, w - 56, 60, "Replace with repo link, QR code, or release note.", 14, TEXT, False)


def add_process_line(slide, x1, y1, x2, y2, color=CYAN):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x1), px(y1), px(x2), px(y2))
    line.line.color.rgb = color
    line.line.width = Pt(2)
    line.line.transparency = 0.25
    return line


def build_deck(output_path, deck_title, deck_subtitle, brand_label, repo_url):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_orb(slide, 120, 120, 640, 640, ORANGE, 0.87)
    add_orb(slide, 1260, 120, 520, 520, PURPLE, 0.88)
    add_orb(slide, 760, 700, 420, 220, CYAN, 0.9)
    add_textbox(slide, 420, 242, 1080, 40, brand_label, 18, PURPLE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 260, 320, 1400, 180, deck_title, 64, WHITE, True, PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 360, 530, 1200, 80, deck_subtitle, 18, TEXT, False, PP_ALIGN.CENTER)
    add_chip(slide, 760, 654, 170, 42, "dark-tech", CYAN)
    add_chip(slide, 950, 654, 170, 42, "source-inspired", ORANGE)
    add_chip(slide, 1140, 654, 170, 42, "editable", GREEN)
    add_footer(slide, "Cover / Key Concept")

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_orb(slide, 1380, 220, 300, 300, PURPLE, 0.9)
    add_header(slide, "01  Evidence / Quote")
    add_textbox(slide, 50, 132, 1820, 86, "Use one large dark card to hold proof, screenshots, or citations.", 28, WHITE, True, PP_ALIGN.CENTER)
    add_quote_card(slide, 250, 280, 1420, 560)
    add_footer(slide, "Evidence / External Validation")

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_orb(slide, 120, 160, 360, 360, ORANGE, 0.9)
    add_header(slide, "02  Comparison")
    add_textbox(slide, 50, 132, 1820, 86, "Hard contrast works better than long explanation.", 28, WHITE, True, PP_ALIGN.CENTER)
    add_comparison_card(slide, 50, 280, 880, 620, "Visual Model", ORANGE, [
        "Can read frames but misses timing context",
        "Often describes what is visible",
        "Weak at narration plus structure coupling",
        "Useful for static screenshots and UI diffs",
    ])
    add_comparison_card(slide, 990, 280, 880, 620, "Omni Model", PURPLE, [
        "Tracks sequence and transitions",
        "Connects audio, visual, and intent",
        "Can infer presentation logic",
        "Better for video reverse-engineering workflows",
    ])
    add_textbox(slide, 930, 552, 60, 60, "VS", 20, CYAN, True, PP_ALIGN.CENTER)
    add_footer(slide, "Comparison / Two Column")

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_orb(slide, 1450, 660, 260, 260, GREEN, 0.9)
    add_header(slide, "03  Leaderboard / Data")
    add_textbox(slide, 50, 132, 1820, 86, "Charts should still feel like part of the same dark system.", 28, WHITE, True, PP_ALIGN.CENTER)
    add_card(slide, 190, 280, 1540, 560, line_color=GREEN, fill_color=CARD_FILL, transparency=0.03)
    add_textbox(slide, 240, 316, 500, 30, "Representative benchmark view", 18, GREEN, True)
    add_bar(slide, 300, 430, 900, 36, "Mimo-V2-Omni", "92", GREEN, muted=False)
    add_bar(slide, 300, 500, 760, 36, "Qwen-Omni", "78", CYAN, muted=True)
    add_bar(slide, 300, 570, 620, 36, "Gemini Flash", "64", ORANGE, muted=True)
    add_bar(slide, 300, 640, 560, 36, "Static vision baseline", "58", RED, muted=True)
    add_textbox(slide, 1240, 430, 360, 180, "Highlight one winner.\nDo not overdecorate the chart.\nThe card itself should carry the style.", 18, TEXT, False)
    add_chip(slide, 1260, 650, 210, 42, "single highlight", GREEN)
    add_footer(slide, "Data / Leaderboard")

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_orb(slide, 70, 740, 320, 220, CYAN, 0.92)
    add_header(slide, "04  Code / CLI")
    add_textbox(slide, 50, 132, 1820, 86, "Pair the command surface with one explanation column.", 28, WHITE, True, PP_ALIGN.CENTER)
    add_code_window(slide, 50, 280, 1120, 660)
    add_card(slide, 1210, 280, 660, 660, line_color=ORANGE, fill_color=CARD_FILL, transparency=0.03)
    add_textbox(slide, 1240, 308, 600, 34, "Narration Notes", 20, ORANGE, True)
    add_textbox(slide, 1240, 368, 590, 420, "Use this column for:\n\n- what the command proves\n- what artifact it outputs\n- why the step exists\n- how to modify the flow next time", 18, TEXT, False)
    add_chip(slide, 1240, 820, 150, 40, "cli", CYAN)
    add_chip(slide, 1410, 820, 150, 40, "analysis", PURPLE)
    add_process_line(slide, 1170, 610, 1210, 610, ORANGE)
    add_footer(slide, "Code / Command Surface")

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_orb(slide, 220, 120, 520, 520, PURPLE, 0.88)
    add_orb(slide, 1240, 520, 430, 430, CYAN, 0.88)
    add_header(slide, "05  Workflow")
    add_textbox(slide, 50, 132, 1820, 86, "The video workflow itself can become a slide.", 28, WHITE, True, PP_ALIGN.CENTER)
    steps = [
        ("Collect", "share link or local mp4", CYAN),
        ("Understand", "Mimo multimodal analysis", PURPLE),
        ("Extract", "structure, evidence, style", ORANGE),
        ("Convert", "deck, report, storyboard", GREEN),
    ]
    x = 150
    for idx, (title, body, color) in enumerate(steps, 1):
        add_card(slide, x, 340, 300, 190, line_color=color, fill_color=CARD_FILL, transparency=0.03)
        add_textbox(slide, x + 24, 364, 252, 30, f"{idx:02d}  {title}", 20, color, True)
        add_textbox(slide, x + 24, 418, 252, 74, body, 16, TEXT, False)
        if idx < len(steps):
            add_process_line(slide, x + 300, 435, x + 360, 435, CYAN)
        x += 360
    add_card(slide, 250, 650, 1420, 150, line_color=CYAN, fill_color=CARD_FILL_2, transparency=0.02)
    add_textbox(slide, 290, 686, 1340, 60, "Conclusion: use one reusable workflow and swap only the prompt plus the output artifact.", 22, WHITE, True, PP_ALIGN.CENTER)
    add_footer(slide, "Workflow / From Video To Assets")

    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_orb(slide, 260, 150, 480, 480, ORANGE, 0.9)
    add_orb(slide, 1320, 540, 360, 360, PURPLE, 0.9)
    add_textbox(slide, 260, 220, 1400, 120, "Build The Next Deck From This Master", 44, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 400, 350, 1120, 50, "One deck family. Multiple AI demos. Same visual system.", 18, TEXT, False, PP_ALIGN.CENTER)
    add_card(slide, 460, 450, 1000, 180, line_color=ORANGE, fill_color=CARD_FILL, transparency=0.03)
    add_textbox(slide, 520, 492, 880, 34, "Countdown / CTA", 20, ORANGE, True, PP_ALIGN.CENTER)
    add_textbox(slide, 520, 540, 880, 46, "3 days 15 hours 59 mins", 28, WHITE, True, PP_ALIGN.CENTER)
    add_repo_card(slide, 610, 700, 700, 190, repo_url)
    add_footer(slide, "Closing / CTA")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def parse_args():
    parser = argparse.ArgumentParser(description="Generate a dark-tech editable PPTX skeleton.")
    parser.add_argument("--output", required=True, help="Absolute or relative output .pptx path")
    parser.add_argument("--deck-title", default="Hunter Alpha", help="Cover title")
    parser.add_argument(
        "--deck-subtitle",
        default="Mimo-V2-Omni source-inspired dark tech presentation skeleton",
        help="Cover subtitle",
    )
    parser.add_argument("--brand-label", default="MIMO-V2-OMNI", help="Small label above title")
    parser.add_argument("--repo-url", default="github.com/your/repo", help="CTA repo or URL text")
    return parser.parse_args()


def main():
    args = parse_args()
    build_deck(
        output_path=Path(args.output).expanduser().resolve(),
        deck_title=args.deck_title,
        deck_subtitle=args.deck_subtitle,
        brand_label=args.brand_label,
        repo_url=args.repo_url,
    )


if __name__ == "__main__":
    main()
